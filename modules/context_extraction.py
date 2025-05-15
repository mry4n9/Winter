import streamlit as st
import requests
from bs4 import BeautifulSoup
import pdfplumber
from pptx import Presentation
from io import BytesIO
from .utils import add_http_if_missing

@st.cache_data(show_spinner=False)
def extract_text_from_url(url: str) -> str:
    """Extracts all text content from a given URL."""
    if not url:
        return ""
    
    processed_url = add_http_if_missing(url)
    try:
        response = requests.get(processed_url, timeout=10)
        response.raise_for_status()  # Raise an exception for HTTP errors
        soup = BeautifulSoup(response.content, 'html.parser')
        
        # Remove script and style elements
        for script_or_style in soup(["script", "style"]):
            script_or_style.decompose()
        
        # Get text
        text = soup.get_text(separator=' ', strip=True)
        return text
    except requests.exceptions.RequestException as e:
        st.error(f"Error fetching URL ({processed_url}): {e}")
        return ""
    except Exception as e:
        st.error(f"Error parsing URL content ({processed_url}): {e}")
        return ""

@st.cache_data(show_spinner=False)
def extract_text_from_pdf(uploaded_file: BytesIO) -> str:
    """Extracts text from an uploaded PDF file."""
    if not uploaded_file:
        return ""
    try:
        with pdfplumber.open(uploaded_file) as pdf:
            text = ""
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        return text.strip()
    except Exception as e:
        st.error(f"Error reading PDF file: {e}")
        return ""

@st.cache_data(show_spinner=False)
def extract_text_from_ppt(uploaded_file: BytesIO) -> str:
    """Extracts text from an uploaded PPTX file."""
    if not uploaded_file:
        return ""
    try:
        prs = Presentation(uploaded_file)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text).strip()
    except Exception as e:
        st.error(f"Error reading PPTX file: {e}")
        return ""

def extract_text_from_uploaded_file(uploaded_file) -> str:
    """
    Detects file type (PDF/PPTX) from an UploadedFile object and extracts text.
    """
    if uploaded_file is None:
        return ""
    
    file_bytes = BytesIO(uploaded_file.getvalue())
    file_name = uploaded_file.name.lower()

    if file_name.endswith(".pdf"):
        return extract_text_from_pdf(file_bytes)
    elif file_name.endswith(".pptx"):
        return extract_text_from_ppt(file_bytes)
    else:
        st.warning(f"Unsupported file type: {uploaded_file.name}. Please upload PDF or PPTX.")
        return ""